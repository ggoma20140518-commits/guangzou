import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import os
import shutil
import tempfile
import re
import base64
from datetime import datetime

# ─── 앱 디렉토리 (배포 환경 호환) ───
APP_DIR = os.path.dirname(os.path.abspath(__file__)) if "__file__" in dir() else os.getcwd()

# ─── 페이지 설정 ───
st.set_page_config(page_title="코스맥스 광저우 대시보드", layout="wide", page_icon="📊")

# ─── 디자인 시스템 (pigment_dispense PPT 참고) ───
COLORS = {
    "primary": "#111827",
    "accent": "#E61E3D",
    "green": "#059669",
    "blue": "#2563EB",
    "orange": "#D97706",
    "purple": "#7C3AED",
    "gray": "#6B7280",
    "light": "#F3F4F6",
    "dark": "#1F2937",
    "muted": "#9CA3AF",
}
CHANNEL_COLORS = {"온라인": "#2563EB", "왕홍": "#7C3AED", "오프라인": "#059669", "수출": "#D97706"}

# ─── 커스텀 CSS ───
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap');

/* 전체 폰트 */
html, body, [class*="css"] { font-family: 'Noto Sans KR', sans-serif; }

/* 사이드바 */
[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #111827 0%, #1F2937 100%);
}
[data-testid="stSidebar"] * { color: #F3F4F6 !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] label p { color: #F3F4F6 !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] section { background: #374151 !important; border-color: #6B7280 !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] section * { color: #D1D5DB !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] button { background: #4B5563 !important; color: #F3F4F6 !important; }
[data-testid="stSidebar"] .stRadio label { font-size: 0.95rem; padding: 6px 0; }

/* 메인 타이틀 */
h1 { color: #111827 !important; font-weight: 900 !important; letter-spacing: -0.5px; border-bottom: 3px solid #E61E3D; padding-bottom: 12px; }
h2, h3 { color: #1F2937 !important; font-weight: 700 !important; }

/* KPI 메트릭 카드 */
[data-testid="stMetric"] {
    background: linear-gradient(135deg, #F9FAFB 0%, #F3F4F6 100%);
    border: 1px solid #E5E7EB;
    border-left: 4px solid #E61E3D;
    border-radius: 12px;
    padding: 16px 20px;
    box-shadow: 0 1px 3px rgba(0,0,0,0.06);
}
[data-testid="stMetric"] label { color: #6B7280 !important; font-size: 0.8rem !important; font-weight: 500 !important; text-transform: uppercase; letter-spacing: 0.5px; }
[data-testid="stMetricValue"] { color: #111827 !important; font-weight: 900 !important; font-size: 1.6rem !important; }
[data-testid="stMetricDelta"] { font-weight: 600 !important; }

/* 데이터프레임 */
[data-testid="stDataFrame"] {
    border: 1px solid #E5E7EB;
    border-radius: 8px;
    overflow: hidden;
}

/* 다운로드 버튼 */
.stDownloadButton > button {
    background: #1F2937 !important;
    color: #F3F4F6 !important;
    border: none !important;
    border-radius: 8px !important;
    font-size: 0.78rem !important;
    font-weight: 500 !important;
    padding: 6px 16px !important;
    transition: all 0.2s;
}
.stDownloadButton > button:hover {
    background: #E61E3D !important;
    transform: translateY(-1px);
    box-shadow: 0 4px 12px rgba(230,30,61,0.3);
}

/* 일반 버튼 */
.stButton > button[kind="primary"] {
    background: #E61E3D !important;
    border: none !important;
    border-radius: 8px !important;
    font-weight: 700 !important;
}

/* 구분선 */
hr { border-color: #E5E7EB !important; margin: 1.5rem 0 !important; }

/* 서브헤더 스타일 */
.stSubheader, h3 {
    position: relative;
    padding-left: 12px !important;
}
</style>
""", unsafe_allow_html=True)

# ─── 데이터 로드 ───
@st.cache_data(show_spinner="데이터 로딩 중...")
def load_data(file, _file_hash=None):
    xls = pd.ExcelFile(file)

    # --- 시트1: 법인별 매출 ---
    df_raw = pd.read_excel(xls, sheet_name=xls.sheet_names[0], header=None)
    months = ["1월", "2월", "3월", "1Q", "4월", "5월", "6월", "2Q", "7월", "8월", "9월", "3Q", "10월", "11월", "12월", "4Q", "TOTAL"]

    entity_data = {}
    for idx, row in df_raw.iterrows():
        cell0 = str(row.iloc[0]).strip() if pd.notna(row.iloc[0]) else ""
        cell1 = str(row.iloc[1]).strip() if pd.notna(row.iloc[1]) else ""
        if cell0 in ["TOTAL", "GZ", "YSGB"] or cell0 == "" and cell1 != "":
            pass  # will parse below

    # 법인별 데이터 파싱 (행 2~6: TOTAL, 7~11: GZ, 12~16: YSGB)
    def parse_entity_block(start_row):
        result = {}
        labels = ["25년매출", "사업계획", "예상실적", "달성률", "성장률"]
        for i, label in enumerate(labels):
            if start_row + i >= len(df_raw):
                result[label] = [0] * 17
                continue
            row = df_raw.iloc[start_row + i]
            values = []
            for j in range(2, 19):  # 열 2~18 (1월~TOTAL)
                v = row.iloc[j]
                try:
                    fv = float(v)
                    values.append(fv if not pd.isna(fv) else 0)
                except (ValueError, TypeError):
                    values.append(0)
            result[label] = values
        return result

    entity_data["TOTAL"] = parse_entity_block(2)
    entity_data["GZ"] = parse_entity_block(7)
    entity_data["YSGB"] = parse_entity_block(12)

    # --- 3번째 시트: 고객사별 정리 데이터 (선택) ---
    df_customers = None
    if len(xls.sheet_names) >= 3:
        try:
            df_customers = pd.read_excel(xls, sheet_name=xls.sheet_names[2])
        except Exception:
            pass
    if df_customers is None:
        df_customers = pd.DataFrame()
    # 숫자 컬럼 변환
    for col in ["25년매출", "25년매출(동기)", "사업계획", "예상매출", "25년매출.1", "사업계획.1", "25년매출(동기).1", "예상매출.1", "사업계획비달성률", "점유율"]:
        if col in df_customers.columns:
            df_customers[col] = pd.to_numeric(df_customers[col], errors="coerce")

    # --- 시트2: 채널별 월별 매출 ---
    df_ch_raw = pd.read_excel(xls, sheet_name=xls.sheet_names[1], header=None)
    channel_monthly = {}
    # 행 3~6: 온라인, 왕홍, 오프라인, 수출 / col14=채널명, col16~27=월별 데이터
    for idx in range(3, 7):
        row = df_ch_raw.iloc[idx]
        ch_name = str(row.iloc[14]).strip() if pd.notna(row.iloc[14]) else ""
        if ch_name in ["온라인", "왕홍", "오프라인", "수출"]:
            values = []
            for j in range(16, 28):  # 12개월
                v = row.iloc[j] if j < len(row) else 0
                try:
                    values.append(float(v))
                except (ValueError, TypeError):
                    values.append(0)
            channel_monthly[ch_name] = values

    # --- 시트2: 채널별 1Q 합계 (전체 고객사 기반) ---
    ch_map_cn = {"线上": "온라인", "红人": "왕홍", "线下": "오프라인", "出口": "수출"}
    channel_1q = {"온라인": {"sales": 0, "plan": 0, "prev": 0},
                  "왕홍": {"sales": 0, "plan": 0, "prev": 0},
                  "오프라인": {"sales": 0, "plan": 0, "prev": 0},
                  "수출": {"sales": 0, "plan": 0, "prev": 0}}
    for idx in range(2, len(df_ch_raw)):
        ch_raw = str(df_ch_raw.iloc[idx, 2]).strip() if pd.notna(df_ch_raw.iloc[idx, 2]) else ""
        ch = ch_map_cn.get(ch_raw, "")
        if not ch:
            continue
        for key, col in [("sales", 8), ("plan", 7), ("prev", 6)]:
            v = df_ch_raw.iloc[idx, col]
            try:
                channel_1q[ch][key] += float(v)
            except (ValueError, TypeError):
                pass

    # --- 시트2: 전체 고객사 테이블 (row 3+가 데이터, row 2가 헤더) ---
    def safe_float(v):
        try:
            f = float(v)
            return f if not pd.isna(f) else 0
        except (ValueError, TypeError):
            return 0

    df_full_customers = []
    for idx in range(3, len(df_ch_raw)):
        cust = df_ch_raw.iloc[idx, 1]
        if pd.isna(cust) or str(cust).strip() == "":
            break  # 빈 행에서 중단
        ch_raw = str(df_ch_raw.iloc[idx, 2]).strip() if pd.notna(df_ch_raw.iloc[idx, 2]) else ""
        ch_kr = ch_map_cn.get(ch_raw, ch_raw if ch_raw else "/")
        new_cust = str(df_ch_raw.iloc[idx, 3]).strip() if pd.notna(df_ch_raw.iloc[idx, 3]) else ""
        tier34 = str(df_ch_raw.iloc[idx, 4]).strip() if pd.notna(df_ch_raw.iloc[idx, 4]) else ""
        rank_change = str(df_ch_raw.iloc[idx, 9]).strip() if pd.notna(df_ch_raw.iloc[idx, 9]) else ""
        prev_rank = str(df_ch_raw.iloc[idx, 10]).strip() if pd.notna(df_ch_raw.iloc[idx, 10]) else ""
        city_tier = str(df_ch_raw.iloc[idx, 12]).strip() if pd.notna(df_ch_raw.iloc[idx, 12]) else ""

        sales_25 = safe_float(df_ch_raw.iloc[idx, 5])
        prev = safe_float(df_ch_raw.iloc[idx, 6])
        plan = safe_float(df_ch_raw.iloc[idx, 7])
        sales = safe_float(df_ch_raw.iloc[idx, 8])

        if sales <= 0 and plan <= 0 and prev <= 0 and sales_25 <= 0:
            continue
        growth = (sales - prev) / abs(prev) if prev else 0
        achieve = sales / plan if plan else 0
        df_full_customers.append({
            "순위": idx - 2, "고객사": cust, "채널": ch_kr, "채널원본": ch_raw,
            "신규처": new_cust, "3,4선": tier34, "도시등급": city_tier,
            "25년매출": sales_25, "25년동기": prev,
            "사업계획": plan, "예상매출": sales,
            "성장률": growth, "달성률": achieve,
            "전년대비순위": rank_change, "25년1Q순위": prev_rank,
        })
    # 시트2 원본 순위 그대로 유지 (스크린샷과 동일)
    df_full_cust = pd.DataFrame(df_full_customers).reset_index(drop=True)
    df_full_cust["순위"] = range(1, len(df_full_cust) + 1)
    # 점유율 계산
    total_sales = df_full_cust["예상매출"].sum()
    df_full_cust["점유율"] = df_full_cust["예상매출"] / total_sales if total_sales else 0

    # --- 채널별 25년 1Q 합계 (월별 데이터 기반) ---
    channel_prev_monthly = {}
    for ch in ["온라인", "왕홍", "오프라인", "수출"]:
        if ch in channel_monthly:
            channel_prev_monthly[ch] = sum(channel_monthly[ch][:3])
        else:
            channel_prev_monthly[ch] = 0

    # --- 채널별 전년대비 비중 (시트2 row 13-16, col 23 = 원본 PPT와 일치) ---
    channel_yoy_ppt = {}
    for idx, ch in [(13, "온라인"), (14, "왕홍"), (15, "오프라인"), (16, "수출")]:
        try:
            channel_yoy_ppt[ch] = float(df_ch_raw.iloc[idx, 23])
        except (ValueError, TypeError):
            channel_yoy_ppt[ch] = 0

    # --- 4번째 시트: 전체 고객사 목록 (선택) ---
    df_all = pd.DataFrame()
    if len(xls.sheet_names) >= 4:
        try:
            df_all = pd.read_excel(xls, sheet_name=xls.sheet_names[3])
            if "25年实绩" in df_all.columns:
                df_all["25年实绩"] = pd.to_numeric(df_all["25年实绩"], errors="coerce")
        except Exception:
            df_all = pd.DataFrame()

    # --- 분기 감지: 시트2 헤더(row3, col I)에서 읽기 → 없으면 자동 감지 ---
    quarter_map = {3: "1Q", 7: "2Q", 11: "3Q", 15: "4Q"}
    quarter_idx_map = {"1Q": 3, "2Q": 7, "3Q": 11, "4Q": 15}
    current_quarter = None
    current_q_idx = 3

    # 시트2 헤더에서 분기 읽기 (예: "1Q 销售 매출액", "2Q매출액", "3Q" 등)
    try:
        header_text = str(df_ch_raw.iloc[2, 8]).strip()
        match = re.search(r'([1-4]Q)', header_text, re.IGNORECASE)
        if match:
            current_quarter = match.group(1).upper()
            current_q_idx = quarter_idx_map.get(current_quarter, 3)
    except Exception:
        pass

    # 헤더에서 못 읽으면 TOTAL 예상실적 기준 자동 감지
    if not current_quarter:
        current_quarter = "1Q"
        for q_idx in [15, 11, 7, 3]:
            val = entity_data["TOTAL"]["예상실적"][q_idx]
            if val and not pd.isna(val) and val > 0:
                current_quarter = quarter_map[q_idx]
                current_q_idx = q_idx
                break

    return entity_data, df_customers, channel_monthly, df_all, months, channel_1q, df_full_cust, channel_prev_monthly, channel_yoy_ppt, current_quarter, current_q_idx


# ─── 파일 업로드 또는 기본 파일 ───
uploaded = st.sidebar.file_uploader("📂 엑셀 파일 업로드 (동일 양식)", type=["xlsx"])

if uploaded:
    # 업로드된 파일을 임시 저장
    tmp_upload = os.path.join(tempfile.gettempdir(), f"cosmax_{uploaded.name}")
    with open(tmp_upload, "wb") as f:
        f.write(uploaded.getvalue())
    data_source = tmp_upload
    file_hash = hash(uploaded.getvalue())  # 파일 내용 기반 캐시 키
    st.sidebar.success(f"✅ {uploaded.name} 로드됨")
else:
    # 기본 파일 자동 탐색 (양식/템플릿 파일 제외)
    app_dir = APP_DIR
    skip_keywords = ["양식", "템플릿", "template", "가이드"]
    xlsx_files = [f for f in os.listdir(app_dir)
                  if f.endswith(".xlsx") and not f.startswith("~")
                  and not any(kw in f for kw in skip_keywords)]
    if xlsx_files:
        default_file = os.path.join(app_dir, xlsx_files[0])
        tmp_data = os.path.join(tempfile.gettempdir(), "cosmax_data.xlsx")
        try:
            shutil.copy2(default_file, tmp_data)
            data_source = tmp_data
        except Exception:
            data_source = default_file
        st.sidebar.caption(f"기본 파일: {xlsx_files[0]}")
        file_hash = os.path.getmtime(default_file)
    else:
        st.warning("📂 동일 양식의 엑셀 파일을 업로드해주세요.")
        st.stop()

entity_data, df_customers, channel_monthly, df_all, months, channel_1q, df_full_cust, channel_prev_monthly, channel_yoy_ppt, current_quarter, current_q_idx = load_data(data_source, _file_hash=file_hash)

# ─── 잔망루피 이미지 로드 ───
loopy_b64 = ""
for ext in ["png", "jpg", "jpeg", "gif", "webp"]:
    lp = os.path.join(APP_DIR, f"loopy.{ext}")
    if os.path.exists(lp):
        with open(lp, "rb") as f:
            loopy_b64 = base64.b64encode(f.read()).decode()
        break

# ─── 모든 페이지 공통 상단 헤더 ───
loopy_img = f'<img src="data:image/png;base64,{loopy_b64}" style="width:42px;height:42px;border-radius:50%;object-fit:cover;border:2px solid #E61E3D;">' if loopy_b64 else ""
st.markdown(f"""
<div style="display:flex;align-items:center;justify-content:space-between;padding:12px 0 8px 0;border-bottom:2px solid #E5E7EB;margin-bottom:16px;">
    <div style="display:flex;align-items:center;gap:12px;">
        {loopy_img}
        <div>
            <span style="font-size:1.4rem;font-weight:900;color:#111827;letter-spacing:-0.5px;">코스맥스 광저우 매출 분석</span>
            <span style="font-size:0.85rem;color:#9CA3AF;margin-left:8px;">for Taehee Kwon</span>
        </div>
    </div>
    <div style="font-size:0.75rem;color:#9CA3AF;">{current_quarter} 업무보고</div>
</div>
""", unsafe_allow_html=True)

# ─── 사이드바 ───
loopy_sidebar = f'<img src="data:image/png;base64,{loopy_b64}" style="width:30px;height:30px;border-radius:50%;object-fit:cover;">' if loopy_b64 else ""
st.sidebar.markdown(f"""
<div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;">
    {loopy_sidebar}
    <div>
        <div style="font-size:0.95rem;font-weight:800;color:#F3F4F6;">코스맥스 광저우 매출 분석</div>
        <div style="font-size:0.7rem;color:#9CA3AF;">for Taehee Kwon</div>
    </div>
</div>
""", unsafe_allow_html=True)
st.sidebar.markdown("---")
page = st.sidebar.radio("📋 메뉴", ["분기별 매출 요약", "고객사별 실적", "채널별 분석", "PPT 다운로드"])

st.sidebar.markdown("---")
st.sidebar.markdown(f"""
<div style="text-align:center;color:#9CA3AF;line-height:1.6;">
    <div style="font-size:0.8rem;">Last updated: {datetime.now().strftime("%Y-%m-%d")}</div>
    <div style="font-size:0.8rem;">Created by <b>[SUJUNG CHOI]</b></div>
</div>
""", unsafe_allow_html=True)

# ─── 공통 헬퍼 ───
def to_excel(df):
    buf = BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False)
    return buf.getvalue()

def growth_pct(cur, prev):
    if prev and prev != 0:
        return (cur - prev) / abs(prev) * 100
    return 0

CHART_TEMPLATE = dict(
    font=dict(family="Noto Sans KR, sans-serif", color="#374151"),
    plot_bgcolor="#FFFFFF",
    paper_bgcolor="#FFFFFF",
    title_font=dict(color="#111827", size=16),
    legend=dict(font=dict(size=11)),
    margin=dict(l=50, r=30, t=40, b=50),
)

# Plotly 차트 설정: 카메라 아이콘으로 PNG 다운로드 가능
PLOTLY_CONFIG = {"toImageButtonOptions": {"format": "png", "scale": 2, "filename": "chart"}, "displayModeBar": True}

def style_fig(fig, h=400):
    fig.update_layout(**CHART_TEMPLATE, height=h)
    fig.update_xaxes(gridcolor="#F3F4F6", linecolor="#E5E7EB")
    fig.update_yaxes(gridcolor="#F3F4F6", linecolor="#E5E7EB")
    return fig

def show_chart(fig, key=""):
    """차트 표시 (Plotly 툴바에서 PNG 저장 가능)"""
    st.plotly_chart(fig, use_container_width=True, config=PLOTLY_CONFIG, key=key if key else None)

def chart_downloads(fig_name, df_data=None, df_name=None, key_prefix=""):
    """데이터 엑셀 다운로드 버튼 (PNG는 차트 툴바 카메라 아이콘 사용)"""
    if df_data is not None:
        st.download_button("📊 엑셀 다운로드", to_excel(df_data), f"{df_name or fig_name}.xlsx", key=f"{key_prefix}_xlsx")

# ═══════════════════════════════════════
# 페이지 1: 1분기 매출 요약
# ═══════════════════════════════════════
if page == "분기별 매출 요약":
    st.title("분기별 예상 매출")

    total_1q = entity_data["TOTAL"]["예상실적"][3]
    gz_1q = entity_data["GZ"]["예상실적"][3]
    ysgb_1q = entity_data["YSGB"]["예상실적"][3]
    total_1q_prev = entity_data["TOTAL"]["25년매출"][3]
    gz_1q_prev = entity_data["GZ"]["25년매출"][3]
    ysgb_1q_prev = entity_data["YSGB"]["25년매출"][3]
    total_1q_plan = entity_data["TOTAL"]["사업계획"][3]
    gz_1q_plan = entity_data["GZ"]["사업계획"][3]
    ysgb_1q_plan = entity_data["YSGB"]["사업계획"][3]

    # KPI 카드
    col1, col2, col3 = st.columns(3)
    col1.metric(f"총 {current_quarter} 예상 매출", f"{total_1q/1000:.0f} 백만 RMB", f"전년비 {growth_pct(total_1q, total_1q_prev):+.0f}%")
    col2.metric(f"GZ 법인 {current_quarter} 매출", f"{gz_1q/1000:.0f} 백만 RMB", f"전년비 {growth_pct(gz_1q, gz_1q_prev):+.0f}%")
    col3.metric(f"YSGB 법인 {current_quarter} 매출", f"{ysgb_1q/1000:.0f} 백만 RMB", f"전년비 {growth_pct(ysgb_1q, ysgb_1q_prev):+.0f}%")

    st.markdown("---")

    # ── 법인별 1Q 비교 차트 (TOTAL/GZ/YSGB) ──
    left, right = st.columns(2)
    with left:
        st.subheader(f"법인별 {current_quarter} 매출 비교")
        entities = ["TOTAL", "GZ", "YSGB"]
        fig = go.Figure()
        fig.add_trace(go.Bar(name="25년 동기", x=entities, y=[entity_data[e]["25년매출"][3]/1000 for e in entities], marker_color="#D1D5DB"))
        fig.add_trace(go.Bar(name="사업계획", x=entities, y=[entity_data[e]["사업계획"][3]/1000 for e in entities], marker_color="#93C5FD"))
        fig.add_trace(go.Bar(name="26년 예상실적", x=entities, y=[entity_data[e]["예상실적"][3]/1000 for e in entities], marker_color="#E61E3D"))
        style_fig(fig)
        fig.update_layout(barmode="group", yaxis_title="백만 RMB")
        show_chart(fig)
        df_1q_compare = pd.DataFrame([{"법인": e, "25년동기": entity_data[e]["25년매출"][3]/1000, "사업계획": entity_data[e]["사업계획"][3]/1000, "26년예상": entity_data[e]["예상실적"][3]/1000} for e in entities])
        chart_downloads("법인별_1Q비교", df_1q_compare, "법인별_1Q비교", "p1_c1")

    with right:
        # 분기에 따라 상반기/하반기 자동 선택
        if current_quarter in ["1Q", "2Q"]:
            half_label = "상반기 (1~6월)"
            month_labels = ["1월", "2월", "3월", "4월", "5월", "6월"]
            # 엑셀 idx: 0=1월,1=2월,2=3월, (3=1Q skip), 4=4월,5=5월,6=6월
            month_indices = [0, 1, 2, 4, 5, 6]
        else:
            half_label = "하반기 (7~12월)"
            month_labels = ["7월", "8월", "9월", "10월", "11월", "12월"]
            # 엑셀 idx: 8=7월,9=8월,10=9월, (11=3Q skip), 12=10월,13=11월,14=12월
            month_indices = [8, 9, 10, 12, 13, 14]

        st.subheader(f"법인별 월별 매출 추이 ({half_label})")

        # 데이터 유효성 확인
        has_data = any(
            entity_data["TOTAL"]["예상실적"][i] > 0
            for i in month_indices
            if not pd.isna(entity_data["TOTAL"]["예상실적"][i])
        )

        if has_data:
            fig2 = go.Figure()
            clr_map = {"TOTAL": "#E61E3D", "GZ": "#2563EB", "YSGB": "#059669"}
            for entity in ["TOTAL", "GZ", "YSGB"]:
                actual = [entity_data[entity]["예상실적"][i]/1000 if entity_data[entity]["예상실적"][i] and not pd.isna(entity_data[entity]["예상실적"][i]) else None for i in month_indices]
                prev = [entity_data[entity]["25년매출"][i]/1000 if entity_data[entity]["25년매출"][i] and not pd.isna(entity_data[entity]["25년매출"][i]) else None for i in month_indices]
                fig2.add_trace(go.Scatter(x=month_labels, y=actual, mode="lines+markers", name=f"{entity} 26년", line=dict(color=clr_map[entity], width=2.5), connectgaps=True))
                fig2.add_trace(go.Scatter(x=month_labels, y=prev, mode="lines", name=f"{entity} 25년", line=dict(color=clr_map[entity], dash="dot", width=1.5), opacity=0.4, connectgaps=True))
            style_fig(fig2)
            fig2.update_layout(yaxis_title="백만 RMB")
            show_chart(fig2)
            chart_downloads("법인별_월별추이", key_prefix="p1_c2")
        else:
            st.info("📭 자료가 없습니다.")

    # ── 인터랙티브: 법인별 월별 상세 분석 ──
    st.markdown("---")
    st.subheader("🔍 법인별 월별 상세 분석")
    ic1, ic2 = st.columns([1, 3])
    with ic1:
        sel_entity = st.selectbox("법인 선택", ["TOTAL", "GZ", "YSGB"], key="i_entity")
        sel_metric = st.radio("지표 선택", ["매출 비교", "달성률 추이", "성장률 추이"], key="i_metric")
        sel_half = st.radio("기간", ["상반기 (1~6월)", "하반기 (7~12월)"], index=0 if current_quarter in ["1Q","2Q"] else 1, key="i_half")

    with ic2:
        if sel_half == "상반기 (1~6월)":
            m_labels = ["1월", "2월", "3월", "4월", "5월", "6월"]
            m_indices = [0, 1, 2, 4, 5, 6]
        else:
            m_labels = ["7월", "8월", "9월", "10월", "11월", "12월"]
            m_indices = [8, 9, 10, 12, 13, 14]

        # 데이터 유효성 확인
        has_interactive_data = any(
            entity_data[sel_entity]["예상실적"][i] > 0
            for i in m_indices
            if not pd.isna(entity_data[sel_entity]["예상실적"][i])
        )

        if not has_interactive_data:
            st.info("📭 자료가 없습니다.")
        elif sel_metric == "매출 비교":
            fig_i = go.Figure()
            actual_m = [entity_data[sel_entity]["예상실적"][i]/1000 if entity_data[sel_entity]["예상실적"][i] and not pd.isna(entity_data[sel_entity]["예상실적"][i]) else 0 for i in m_indices]
            plan_m = [entity_data[sel_entity]["사업계획"][i]/1000 for i in m_indices]
            prev_m = [entity_data[sel_entity]["25년매출"][i]/1000 for i in m_indices]
            fig_i.add_trace(go.Bar(name="25년 동기", x=m_labels, y=prev_m, marker_color="#D1D5DB"))
            fig_i.add_trace(go.Bar(name="사업계획", x=m_labels, y=plan_m, marker_color="#93C5FD"))
            fig_i.add_trace(go.Bar(name="26년 실적", x=m_labels, y=actual_m, marker_color="#E61E3D"))
            style_fig(fig_i, 380)
            fig_i.update_layout(barmode="group", yaxis_title="백만 RMB", title=f"{sel_entity} 월별 매출 비교")
            show_chart(fig_i)
        elif sel_metric == "달성률 추이":
            rates = []
            for i in m_indices:
                p = entity_data[sel_entity]["사업계획"][i]
                a = entity_data[sel_entity]["예상실적"][i]
                rates.append(a / p * 100 if p and a and not pd.isna(a) else 0)
            fig_i = go.Figure()
            fig_i.add_trace(go.Scatter(x=m_labels, y=rates, mode="lines+markers+text", text=[f"{r:.0f}%" if r else "" for r in rates], textposition="top center", line=dict(color="#E61E3D", width=3), marker=dict(size=10)))
            fig_i.add_hline(y=100, line_dash="dash", line_color="#059669", annotation_text="100% 달성")
            style_fig(fig_i, 380)
            fig_i.update_layout(yaxis_title="달성률 (%)", title=f"{sel_entity} 월별 사업계획 달성률")
            show_chart(fig_i)
        else:
            growths = []
            for i in m_indices:
                a = entity_data[sel_entity]["예상실적"][i]
                p = entity_data[sel_entity]["25년매출"][i]
                growths.append(growth_pct(a, p) if p and a and not pd.isna(a) else 0)
            colors = ["#059669" if g >= 0 else "#E61E3D" for g in growths]
            fig_i = go.Figure()
            fig_i.add_trace(go.Bar(x=m_labels, y=growths, marker_color=colors, text=[f"{g:+.0f}%" if g else "" for g in growths], textposition="outside"))
            fig_i.add_hline(y=0, line_color="#9CA3AF")
            style_fig(fig_i, 380)
            fig_i.update_layout(yaxis_title="전년 대비 성장률 (%)", title=f"{sel_entity} 월별 성장률")
            show_chart(fig_i)

    st.markdown("---")

    # ── 달성률 & 성장률 테이블 ──
    st.subheader("법인별 사업계획 달성률 및 성장률")
    summary_data = []
    for entity in ["TOTAL", "GZ", "YSGB"]:
        row = {"법인": entity}
        for i, m in enumerate(["1월", "2월", "3월", "1Q", "4월", "5월", "6월", "2Q"]):
            plan = entity_data[entity]["사업계획"][i]
            actual = entity_data[entity]["예상실적"][i]
            prev = entity_data[entity]["25년매출"][i]
            rate = actual / plan * 100 if plan else 0
            growth = growth_pct(actual, prev) if prev else 0
            row[f"{m} 달성률"] = f"{rate:.0f}%" if actual else "-"
            row[f"{m} 성장률"] = f"{growth:+.0f}%" if actual else "-"
        summary_data.append(row)
    df_summary = pd.DataFrame(summary_data)
    st.dataframe(df_summary, use_container_width=True, hide_index=True)
    st.download_button("달성률/성장률 엑셀 다운로드", to_excel(df_summary), "법인별_달성률_성장률.xlsx", key="dl_summary")

    # ── 법인별 연간 실적 테이블 (단위: 백만 RMB) ──
    st.subheader("법인별 연간 실적 상세 (단위: 백만 RMB)")
    month_cols = ["1월", "2월", "3월", "1Q", "4월", "5월", "6월", "2Q", "7월", "8월", "9월", "3Q", "10월", "11월", "12월", "4Q", "TOTAL"]
    annual_rows = []
    for entity in ["TOTAL", "GZ", "YSGB"]:
        for label in ["25년매출", "사업계획", "예상실적"]:
            row = {"법인": entity, "구분": label}
            for mi, m in enumerate(month_cols):
                v = entity_data[entity][label][mi]
                row[m] = round(v / 1000) if v and not pd.isna(v) else 0
            annual_rows.append(row)
    df_annual = pd.DataFrame(annual_rows)
    st.dataframe(df_annual, use_container_width=True, hide_index=True)
    st.download_button("연간 실적 엑셀 다운로드", to_excel(df_annual), "법인별_연간실적.xlsx", key="dl_annual")


# ═══════════════════════════════════════
# 페이지 2: 고객사별 실적
# ═══════════════════════════════════════
elif page == "고객사별 실적":
    st.title(f"{current_quarter} TOP 고객사 실적")

    top_n = st.slider("상위 고객사 수", 5, 50, 30)
    df_top = df_full_cust.head(top_n).copy()
    gz_1q_total = entity_data["GZ"]["예상실적"][3]

    # KPI (시트2 전체 데이터 기반, 반올림 합산 = PPT/표와 일치)
    top10_sales_r = sum(round(v / 1000) for v in df_full_cust.head(10)["예상매출"])
    top30_sales_r = sum(round(v / 1000) for v in df_full_cust.head(30)["예상매출"])
    top10_raw = df_full_cust.head(10)["예상매출"].sum()
    top30_raw = df_full_cust.head(30)["예상매출"].sum()
    col1, col2, col3 = st.columns(3)
    col1.metric("TOP 10 매출 합계", f"{top10_sales_r} 백만 RMB", f"GZ의 {top10_raw/gz_1q_total*100:.0f}%")
    col2.metric("TOP 30 매출 합계", f"{top30_sales_r} 백만 RMB", f"GZ의 {top30_raw/gz_1q_total*100:.0f}%")
    total_cust = len(df_full_cust[df_full_cust["예상매출"] > 0])
    col3.metric("매출 발생 고객사 수", f"{total_cust}개사")

    st.markdown("---")

    # ── 고객사별 매출 차트 ──
    left, right = st.columns(2)
    with left:
        st.subheader(f"TOP {top_n} 고객사 예상 매출")
        fig = px.bar(df_top, x="고객사", y="예상매출", color="채널", color_discrete_map=CHANNEL_COLORS)
        style_fig(fig, 450)
        fig.update_layout(yaxis_title="만 RMB", xaxis_tickangle=-45)
        show_chart(fig)
        chart_downloads(f"TOP{top_n}_매출", key_prefix="p2_c1")

    with right:
        st.subheader("고객사별 계획 대비 달성률")
        df_rate = df_top[df_top["달성률"] > 0].copy()
        fig2 = px.bar(df_rate, x="고객사", y="달성률", color_discrete_sequence=[COLORS["accent"]])
        fig2.add_hline(y=0.25, line_dash="dash", line_color="#059669", annotation_text="1Q 목표(25%)")
        style_fig(fig2, 450)
        fig2.update_layout(yaxis_title="달성률", yaxis_tickformat=".0%", xaxis_tickangle=-45)
        show_chart(fig2)
        chart_downloads("달성률", key_prefix="p2_c2")

    # ── 매출 점유율 & 전년 비교 ──
    left2, right2 = st.columns(2)
    with left2:
        st.subheader("매출 점유율 (TOP 10)")
        top10 = df_full_cust.head(10).copy()
        fig3 = px.pie(top10, names="고객사", values="예상매출", color_discrete_sequence=["#E61E3D","#2563EB","#059669","#D97706","#7C3AED","#EC4899","#14B8A6","#F59E0B","#6366F1","#8B5CF6"])
        fig3.update_traces(textposition="inside", textinfo="label+percent")
        style_fig(fig3, 400)
        show_chart(fig3)
        chart_downloads("매출점유율_TOP10", key_prefix="p2_c3")

    with right2:
        st.subheader("전년 동기 대비 매출 변화 (TOP 15)")
        top15 = df_full_cust.head(15).copy()
        top15["전년대비"] = top15["예상매출"] - top15["25년동기"]
        top15["color"] = top15["전년대비"].apply(lambda x: "증가" if x > 0 else "감소")
        fig4 = px.bar(top15, x="고객사", y="전년대비", color="color",
                      color_discrete_map={"증가": "#059669", "감소": "#E61E3D"})
        style_fig(fig4, 400)
        fig4.update_layout(yaxis_title="만 RMB (전년 대비 증감)", xaxis_tickangle=-45, showlegend=False)
        show_chart(fig4)
        chart_downloads("전년대비_증감", top15[["고객사","예상매출","25년동기","전년대비"]], "전년대비_증감", "p2_c4")

    # ── 인터랙티브: 고객사 매출 vs 성장률 산점도 ──
    st.markdown("---")
    st.subheader("🔍 고객사 매출 vs 성장률 분석")
    scatter_data = df_full_cust[(df_full_cust["예상매출"] > 0) & (df_full_cust["25년동기"] > 0)].head(40).copy()
    scatter_data["매출(백만)"] = scatter_data["예상매출"] / 1000
    scatter_data["성장률(%)"] = scatter_data["성장률"] * 100
    scatter_data["사업계획(백만)"] = scatter_data["사업계획"] / 1000

    sc1, sc2 = st.columns([1, 3])
    with sc1:
        x_axis = st.selectbox("X축", ["매출(백만)", "사업계획(백만)"], key="sc_x")
        y_axis = st.selectbox("Y축", ["성장률(%)", "달성률"], key="sc_y")
        if y_axis == "달성률":
            scatter_data["달성률(%)"] = scatter_data["달성률"] * 100
            y_col = "달성률(%)"
        else:
            y_col = y_axis
        ch_filter = st.multiselect("채널 필터", ["온라인", "왕홍", "오프라인", "수출"], default=["온라인", "왕홍", "오프라인", "수출"], key="sc_ch")
    with sc2:
        sc_filtered = scatter_data[scatter_data["채널"].isin(ch_filter)]
        fig_sc = px.scatter(sc_filtered, x=x_axis, y=y_col, color="채널", size="매출(백만)",
                            hover_name="고객사", hover_data=["매출(백만)", "성장률(%)", "채널"],
                            color_discrete_map=CHANNEL_COLORS, size_max=40)
        fig_sc.add_hline(y=0, line_dash="dot", line_color="#9CA3AF")
        fig_sc.add_vline(x=sc_filtered[x_axis].median(), line_dash="dot", line_color="#9CA3AF")
        style_fig(fig_sc, 450)
        fig_sc.update_layout(title="고객사 포지셔닝 맵 (버블 크기 = 매출)")
        show_chart(fig_sc)
        chart_downloads("고객사_포지셔닝", key_prefix="p2_sc")

    st.markdown("---")

    # ── 채널별 고객사 분포 ──
    st.subheader("채널별 고객사 수 및 매출 비중")
    ch_summary = df_full_cust[df_full_cust["예상매출"] > 0].groupby("채널").agg(
        고객수=("고객사", "count"),
        매출합계=("예상매출", "sum"),
        평균매출=("예상매출", "mean"),
    ).reset_index()
    ch_summary["매출합계(백만)"] = (ch_summary["매출합계"] / 1000).round(0)
    ch_summary["평균매출(백만)"] = (ch_summary["평균매출"] / 1000).round(1)
    st.dataframe(ch_summary[["채널", "고객수", "매출합계(백만)", "평균매출(백만)"]], use_container_width=True, hide_index=True)

    # ── 상세 테이블 (시트2 원본 데이터 기반, PPT와 동일) ──
    st.subheader(f"TOP {top_n} 고객사 상세 데이터 (단위: 백만 RMB)")
    df_display = df_top.copy()
    df_display["25년매출(백만)"] = df_display["25년매출"].apply(lambda x: round(x / 1000) if x else 0)
    df_display["사업계획(백만)"] = df_display["사업계획"].apply(lambda x: round(x / 1000) if x else 0)
    df_display["25년동기(백만)"] = df_display["25년동기"].apply(lambda x: round(x / 1000) if x else 0)
    df_display["예상매출(백만)"] = df_display["예상매출"].apply(lambda x: round(x / 1000) if x else 0)
    df_display["성장률(%)"] = df_display["성장률"].apply(lambda x: f"{x:.0%}" if x and x != 0 else "-")
    df_display["달성률(%)"] = df_display["달성률"].apply(lambda x: f"{x:.0%}" if x else "-")
    df_display["점유율(%)"] = df_display["점유율"].apply(lambda x: f"{x:.0%}" if x else "-")
    show_cols = ["순위", "고객사", "채널", "신규처", "3,4선", "25년매출(백만)", "사업계획(백만)", "25년동기(백만)", "예상매출(백만)", "성장률(%)", "달성률(%)", "점유율(%)", "전년대비순위"]
    st.dataframe(df_display[show_cols], use_container_width=True, hide_index=True)
    st.download_button("고객사 실적 엑셀 다운로드", to_excel(df_display[show_cols]), "고객사별_실적.xlsx", key="dl_cust")


# ═══════════════════════════════════════
# 페이지 3: 채널별 분석
# ═══════════════════════════════════════
elif page == "채널별 분석":
    st.title(f"{current_quarter} 유통 채널별 실적 분석")

    # 시트2 전체 데이터 기반 채널별 집계
    total_ch_sales = sum(channel_1q[ch]["sales"] for ch in channel_1q)

    # KPI 카드 (전년비 = 시트2 col23 "전년대비 비중", 원본 PPT와 일치)
    cols = st.columns(4)
    for i, ch in enumerate(["온라인", "왕홍", "오프라인", "수출"]):
        sales = channel_1q[ch]["sales"]
        share = sales / total_ch_sales * 100 if total_ch_sales else 0
        yoy = channel_yoy_ppt.get(ch, 0)
        arrow = "+" if yoy > 0 else ""
        cols[i].metric(ch, f"{sales/1000:.0f} 백만 RMB", f"전년비 {arrow}{yoy:.0f}% | 비중 {share:.0f}%")

    st.markdown("---")

    # ── 채널별 매출/비중/달성률 ──
    left, right = st.columns(2)
    with left:
        st.subheader("채널별 매출 비중 (도넛)")
        ch_pie = pd.DataFrame([{"채널": ch, "매출": channel_1q[ch]["sales"]} for ch in ["온라인", "왕홍", "오프라인", "수출"]])
        fig = px.pie(ch_pie, names="채널", values="매출", color="채널", color_discrete_map=CHANNEL_COLORS, hole=0.4)
        fig.update_traces(textposition="inside", textinfo="label+percent", textfont_size=13)
        style_fig(fig, 400)
        show_chart(fig)
        chart_downloads("채널별_매출비중", ch_pie, "채널별_매출비중", "p3_c1")

    with right:
        st.subheader("채널별 사업계획 달성률")
        ch_rate_data = []
        for ch in ["온라인", "왕홍", "오프라인", "수출"]:
            s = channel_1q[ch]["sales"]
            p = channel_1q[ch]["plan"]
            ch_rate_data.append({"채널": ch, "달성률": s / p if p else 0})
        df_ch_rate = pd.DataFrame(ch_rate_data)
        fig2 = px.bar(df_ch_rate, x="채널", y="달성률", color="채널", color_discrete_map=CHANNEL_COLORS)
        fig2.add_hline(y=0.25, line_dash="dash", line_color="#E61E3D", annotation_text="1Q 목표(25%)")
        style_fig(fig2, 400)
        fig2.update_layout(yaxis_tickformat=".0%", showlegend=False)
        show_chart(fig2)
        chart_downloads("채널별_달성률", df_ch_rate, "채널별_달성률", "p3_c2")

    # ── 채널별 TOP 5 고객사 ──
    st.subheader("채널별 TOP 5 고객사")
    ch_cols = st.columns(4)
    for i, ch in enumerate(["온라인", "왕홍", "오프라인", "수출"]):
        with ch_cols[i]:
            st.markdown(f"**{ch}**")
            df_ch_top = df_full_cust[df_full_cust["채널"] == ch].head(5)
            for _, r in df_ch_top.iterrows():
                st.markdown(f"- {r['고객사']}: **{r['예상매출']/1000:.0f}**백만")

    # ── 채널별 전년 대비 비교 ──
    st.subheader("채널별 전년 동기 대비 비교")
    ch_compare = pd.DataFrame([
        {"채널": ch,
         "25년 1Q": channel_prev_monthly.get(ch, 0) / 1000,
         "26년 예상": channel_1q[ch]["sales"] / 1000,
         "사업계획": channel_1q[ch]["plan"] / 1000}
        for ch in ["온라인", "왕홍", "오프라인", "수출"]
    ])
    fig3 = go.Figure()
    fig3.add_trace(go.Bar(name="25년 1Q", x=ch_compare["채널"], y=ch_compare["25년 1Q"], marker_color="#D1D5DB"))
    fig3.add_trace(go.Bar(name="26년 사업계획", x=ch_compare["채널"], y=ch_compare["사업계획"], marker_color="#93C5FD"))
    fig3.add_trace(go.Bar(name="26년 예상실적", x=ch_compare["채널"], y=ch_compare["26년 예상"], marker_color="#E61E3D"))
    style_fig(fig3, 400)
    fig3.update_layout(barmode="group", yaxis_title="백만 RMB")
    show_chart(fig3)
    chart_downloads("채널별_전년대비", ch_compare, "채널별_전년대비", "p3_c3")

    # ── 25년 월별 추이 ──
    if channel_monthly:
        st.subheader("25년 채널별 월별 매출 추이")
        month_names = ["1월", "2월", "3월", "4월", "5월", "6월", "7월", "8월", "9월", "10월", "11월", "12월"]
        fig4 = go.Figure()
        for ch, vals in channel_monthly.items():
            display_vals = vals[:12] if len(vals) >= 12 else vals
            fig4.add_trace(go.Scatter(x=month_names[:len(display_vals)], y=display_vals, mode="lines+markers", name=ch, line=dict(color=CHANNEL_COLORS.get(ch, "#333"))))
        style_fig(fig4, 400)
        fig4.update_layout(yaxis_title="만 RMB")
        show_chart(fig4)
        df_monthly_export = pd.DataFrame({ch: vals[:12] for ch, vals in channel_monthly.items()}, index=month_names)
        chart_downloads("채널별_월별추이", df_monthly_export.reset_index().rename(columns={"index": "월"}), "채널별_월별추이", "p3_c4")

    # ── 인터랙티브: 채널 드릴다운 분석 ──
    st.markdown("---")
    st.subheader("🔍 채널별 고객사 드릴다운")
    dd1, dd2 = st.columns([1, 3])
    with dd1:
        sel_ch = st.selectbox("채널 선택", ["온라인", "왕홍", "오프라인", "수출"], key="dd_ch")
        dd_top_n = st.slider("표시 고객수", 5, 30, 15, key="dd_n")
    with dd2:
        df_dd = df_full_cust[df_full_cust["채널"] == sel_ch].head(dd_top_n).copy()
        df_dd["매출(백만)"] = df_dd["예상매출"] / 1000
        df_dd["동기(백만)"] = df_dd["25년동기"] / 1000

        fig_dd = go.Figure()
        fig_dd.add_trace(go.Bar(name="25년 동기", x=df_dd["고객사"], y=df_dd["동기(백만)"], marker_color="#D1D5DB"))
        fig_dd.add_trace(go.Bar(name="26년 예상", x=df_dd["고객사"], y=df_dd["매출(백만)"], marker_color=CHANNEL_COLORS.get(sel_ch, "#E61E3D")))
        style_fig(fig_dd, 400)
        fig_dd.update_layout(barmode="group", title=f"{sel_ch} 채널 TOP {dd_top_n} 고객사", yaxis_title="백만 RMB", xaxis_tickangle=-45)
        show_chart(fig_dd)

    # 선택 채널 요약 KPI
    ch_kpi = st.columns(4)
    ch_s = channel_1q[sel_ch]["sales"]
    ch_p = channel_1q[sel_ch]["plan"]
    ch_prev = channel_prev_monthly.get(sel_ch, 0)
    ch_kpi[0].metric(f"{sel_ch} 매출", f"{ch_s/1000:.0f} 백만")
    ch_kpi[1].metric("달성률", f"{ch_s/ch_p*100:.0f}%" if ch_p else "-")
    ch_kpi[2].metric("전년비", f"{channel_yoy_ppt.get(sel_ch, 0):+.0f}%")
    ch_kpi[3].metric("고객수", f"{len(df_full_cust[(df_full_cust['채널']==sel_ch) & (df_full_cust['예상매출']>0)])}개사")

    st.markdown("---")

    # ── 채널 상세 테이블 ──
    st.subheader("채널별 상세 데이터")
    ch_detail = pd.DataFrame([
        {"채널": ch,
         "예상매출(백만)": round(channel_1q[ch]["sales"] / 1000),
         "사업계획(백만)": round(channel_1q[ch]["plan"] / 1000),
         "25년1Q(백만)": round(channel_prev_monthly.get(ch, 0) / 1000),
         "달성률": f"{channel_1q[ch]['sales']/channel_1q[ch]['plan']*100:.0f}%" if channel_1q[ch]["plan"] else "-",
         "전년비": f"{channel_yoy_ppt.get(ch, 0):+.0f}%",
         "비중": f"{channel_1q[ch]['sales']/total_ch_sales*100:.0f}%" if total_ch_sales else "-",
         "고객수": len(df_full_cust[(df_full_cust["채널"] == ch) & (df_full_cust["예상매출"] > 0)])}
        for ch in ["온라인", "왕홍", "오프라인", "수출"]
    ])
    st.dataframe(ch_detail, use_container_width=True, hide_index=True)
    st.download_button("채널별 분석 엑셀 다운로드", to_excel(ch_detail), "채널별_분석.xlsx", key="dl_channel")


# ═══════════════════════════════════════
# 페이지 4: PPT 다운로드
# ═══════════════════════════════════════
elif page == "PPT 다운로드":
    st.title("분석 결과 PPT 다운로드")
    st.markdown("기존 PPT 양식을 유지하면서 최신 데이터로 업데이트합니다.")

    report_date = st.date_input("보고 일자", value=pd.Timestamp.now())

    # 템플릿 PPT 경로
    template_file = os.path.join(APP_DIR, "코스맥스 광저우 주간 보고_20260401.pptx")
    if not os.path.exists(template_file):
        st.error("템플릿 PPT 파일이 없습니다: 코스맥스 광저우 주간 보고_20260401.pptx")
        st.stop()

    if st.button("PPT 생성", type="primary"):
        with st.spinner("PPT 생성 중..."):

            # ── 헬퍼 함수들 ──
            def find_shape(slide, left, top, tol=50000):
                """위치 기반으로 shape 검색"""
                for s in slide.shapes:
                    if abs(s.left - left) < tol and abs(s.top - top) < tol:
                        return s
                return None

            def set_text_keep_format(shape, new_text, auto_fit=False):
                """텍스트 교체하되 기존 폰트 서식 유지"""
                if not shape or not shape.has_text_frame:
                    return
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        para.runs[0].text = str(new_text)
                        for run in para.runs[1:]:
                            run.text = ""
                        break
                else:
                    shape.text_frame.paragraphs[0].text = str(new_text)
                # 텍스트 자동 축소 (박스 넘침 방지)
                if auto_fit:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
                    if bodyPr is not None:
                        # 기존 autofit 요소 제거
                        for child in bodyPr.findall(qn("a:noAutofit")):
                            bodyPr.remove(child)
                        for child in bodyPr.findall(qn("a:spAutoFit")):
                            bodyPr.remove(child)
                        normAuto = bodyPr.find(qn("a:normAutofit"))
                        if normAuto is None:
                            normAuto = etree.SubElement(bodyPr, qn("a:normAutofit"))
                        normAuto.set("fontScale", "60000")

            def clear_all_text(shape):
                """shape 내 모든 텍스트 제거 (서식 유지)"""
                if not shape or not shape.has_text_frame:
                    return
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

            def set_table_cell(table, row, col, text):
                """테이블 셀 텍스트 교체 (서식 유지)"""
                cell = table.cell(row, col)
                for para in cell.text_frame.paragraphs:
                    if para.runs:
                        para.runs[0].text = str(text)
                        for run in para.runs[1:]:
                            run.text = ""
                        return
                cell.text_frame.paragraphs[0].text = str(text)

            def fig_to_image(fig, width=800, height=400):
                """Plotly 차트를 PNG 이미지로 변환 (kaleido 실패 시 matplotlib 폴백)"""
                try:
                    img_bytes = fig.to_image(format="png", width=width, height=height, scale=2)
                    return BytesIO(img_bytes)
                except Exception:
                    # kaleido/Chrome 없는 환경 (Streamlit Cloud 등) → matplotlib 폴백
                    import matplotlib
                    matplotlib.use("Agg")
                    import matplotlib.pyplot as plt
                    import matplotlib.font_manager as fm
                    buf = BytesIO()
                    fig_mpl = fig.to_dict()
                    mpl_fig, ax = plt.subplots(figsize=(width/100, height/100), dpi=150)
                    # 간단한 바 차트 렌더링
                    traces = fig_mpl.get("data", [])
                    x_all = []
                    for t in traces:
                        if t.get("type") == "bar" and t.get("x") and t.get("y"):
                            x_labels = list(t["x"])
                            x_all = x_labels
                            ax.bar(x_labels, list(t["y"]), label=t.get("name", ""), alpha=0.8)
                        elif t.get("type") == "pie":
                            labels = list(t.get("labels", []))
                            values = list(t.get("values", []))
                            ax.pie(values, labels=labels, autopct="%1.0f%%")
                    if x_all:
                        ax.legend(fontsize=7)
                    ax.set_ylabel(fig_mpl.get("layout", {}).get("yaxis", {}).get("title", {}).get("text", ""))
                    plt.tight_layout()
                    plt.savefig(buf, format="png", bbox_inches="tight")
                    plt.close(mpl_fig)
                    buf.seek(0)
                    return buf

            def growth_pct(cur, prev):
                if prev and prev != 0:
                    return (cur - prev) / abs(prev) * 100
                return 0

            # ── 템플릿 복사 후 열기 ──
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            tmp.close()
            shutil.copy2(template_file, tmp.name)
            prs = Presentation(tmp.name)

            # ════════════════════════════════════
            # 슬라이드 1: 표지 - 날짜만 업데이트
            # ════════════════════════════════════
            slide1 = prs.slides[0]
            date_shape = find_shape(slide1, 481013, 4305300)
            if date_shape:
                set_text_keep_format(date_shape, report_date.strftime("%Y.%m.%d"))

            # ════════════════════════════════════
            # 슬라이드 2: 1분기 매출 요약
            # ════════════════════════════════════
            slide2 = prs.slides[1]

            total_1q = entity_data["TOTAL"]["예상실적"][3]
            gz_1q = entity_data["GZ"]["예상실적"][3]
            ysgb_1q = entity_data["YSGB"]["예상실적"][3]
            total_prev = entity_data["TOTAL"]["25년매출"][3]
            gz_prev = entity_data["GZ"]["25년매출"][3]
            ysgb_prev = entity_data["YSGB"]["25년매출"][3]

            # KPI 큰 숫자 (Text 6) - 3개 법인
            for left, top, val in [
                (578893, 969378, total_1q),
                (4470026, 969378, gz_1q),
                (8332405, 969378, ysgb_1q),
            ]:
                s = find_shape(slide2, left, top)
                if s:
                    set_text_keep_format(s, f"{val/1000:.0f} 백만 RMB", auto_fit=True)

            # 성장률 텍스트 (Text 7) - 3개 법인
            for left, top, cur, prev in [
                (578893, 1788358, total_1q, total_prev),
                (4467883, 1797484, gz_1q, gz_prev),
                (8331983, 1788358, ysgb_1q, ysgb_prev),
            ]:
                s = find_shape(slide2, left, top)
                if s:
                    g = growth_pct(cur, prev)
                    if abs(g) < 0.5:
                        set_text_keep_format(s, "= 전년 동기 대비 0% 성장")
                    else:
                        arrow = "▲" if g > 0 else "▼"
                        set_text_keep_format(s, f"{arrow} 전년 동기 대비 {g:+.0f}% 성장")
                    if s.text_frame.paragraphs[0].runs:
                        run = s.text_frame.paragraphs[0].runs[0]
                        if g > 0:
                            run.font.color.rgb = RGBColor(0x38, 0x57, 0x23)
                        elif g < 0:
                            run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)

            # 이전 보고 대비 → 공란
            for left, top in [(2688255, 1427138), (6563719, 1427138), (10239792, 1431316)]:
                s = find_shape(slide2, left, top)
                if s:
                    set_text_keep_format(s, "")

            # 표 20: 법인별 실적 테이블 (16x13)
            table_shape = None
            for s in slide2.shapes:
                if s.has_table and len(list(s.table.rows)) == 16:
                    table_shape = s
                    break

            if table_shape:
                tbl = table_shape.table
                # PPT 테이블 11개 기간 → 엑셀 17개 기간 인덱스 매핑
                # PPT col: 2=1월,3=2월,4=3월,5=1Q,6=4月,7=5月,8=6月,9=2Q,10=3Q,11=4Q,12=TOTAL
                # 엑셀 idx: 0=1월,1=2월,2=3월,3=1Q,4=4월,5=5월,6=6월,7=2Q,...,11=3Q,...,15=4Q,16=TOTAL
                ppt_to_excel = {2:0, 3:1, 4:2, 5:3, 6:4, 7:5, 8:6, 9:7, 10:11, 11:15, 12:16}

                for ent_idx, entity in enumerate(["TOTAL", "GZ", "YSGB"]):
                    base_row = 1 + ent_idx * 5
                    for li, label in enumerate(["25년매출", "사업계획", "예상실적", "달성률", "성장률"]):
                        row_idx = base_row + li
                        if row_idx >= len(list(tbl.rows)):
                            break
                        for ppt_col, excel_idx in ppt_to_excel.items():
                            if ppt_col >= len(list(tbl.columns)):
                                break
                            val = entity_data[entity][label][excel_idx] if excel_idx < len(entity_data[entity][label]) else 0
                            if label in ["달성률", "성장률"]:
                                if val == 0 or pd.isna(val):
                                    txt = "- "
                                elif label == "성장률" and val == -1:
                                    txt = "- "
                                elif label == "성장률" and val < -0.99:
                                    txt = "- "  # 미래 데이터
                                else:
                                    pct = val * 100 if abs(val) < 2 else val
                                    txt = f"{pct:.0f}%"
                            else:
                                if pd.isna(val) or val == 0:
                                    txt = "- "
                                elif abs(val/1000) >= 1000:
                                    txt = f"{val/1000:,.0f} "
                                else:
                                    txt = f"{val/1000:.0f} "
                            set_table_cell(tbl, row_idx, ppt_col, txt)

            # 차트 53 → 삭제 후 이미지로 교체
            chart_shape = None
            for s in slide2.shapes:
                if s.shape_type == 3:  # CHART
                    chart_shape = s
                    break
            if chart_shape:
                cl, ct, cw, ch_ = chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height
                sp = chart_shape._element
                sp.getparent().remove(sp)

                # 법인별 1Q 매출 비교 차트 (TOTAL, GZ, YSGB)
                entities = ["TOTAL", "GZ", "YSGB"]
                colors_bar = ["#203864", "#2980B9", "#27AE60"]
                vals_1q = [entity_data[e]["예상실적"][3] / 1000 for e in entities]
                prev_1q = [entity_data[e]["25년매출"][3] / 1000 for e in entities]
                plan_1q = [entity_data[e]["사업계획"][3] / 1000 for e in entities]

                fig_chart = go.Figure()
                fig_chart.add_trace(go.Bar(name="25년 동기", x=entities, y=prev_1q,
                                           marker_color="#B0B0B0", opacity=0.6))
                fig_chart.add_trace(go.Bar(name="사업계획", x=entities, y=plan_1q,
                                           marker_color="#A0C4E8", opacity=0.6))
                fig_chart.add_trace(go.Bar(name="26년 예상실적", x=entities, y=vals_1q,
                                           marker_color="#203864"))
                fig_chart.update_layout(
                    barmode="group", yaxis_title="백만 RMB",
                    margin=dict(l=50, r=20, t=20, b=40),
                    font=dict(size=12), plot_bgcolor="white",
                    legend=dict(orientation="h", yanchor="bottom", y=1.02),
                )
                slide2.shapes.add_picture(fig_to_image(fig_chart, 600, 500), cl, ct, cw, ch_)

            # ════════════════════════════════════
            # 슬라이드 3: TOP 30 고객사 실적
            # ════════════════════════════════════
            slide3 = prs.slides[2]

            # 시트2 전체 고객사 기반 TOP 10/30 합계 (반올림 합산 = 표와 일치)
            top10_sales_raw = df_full_cust.head(10)["예상매출"].sum()
            top30_sales_raw = df_full_cust.head(30)["예상매출"].sum()
            top10_sales_rounded = sum(round(v / 1000) for v in df_full_cust.head(10)["예상매출"])
            top30_sales_rounded = sum(round(v / 1000) for v in df_full_cust.head(30)["예상매출"])
            gz_1q_total = entity_data["GZ"]["예상실적"][3]

            # TOP 10 매출 값 (반올림 합산)
            s = find_shape(slide3, 578893, 969378)
            if s:
                set_text_keep_format(s, f"{top10_sales_rounded} 백만 RMB", auto_fit=True)

            # TOP 10 점유율
            s = find_shape(slide3, 597372, 1787352)
            if s:
                set_text_keep_format(s, f"전체 매출의 {top10_sales_raw/gz_1q_total*100:.0f}% ", auto_fit=True)

            # TOP 30 매출 값 (반올림 합산)
            s = find_shape(slide3, 581319, 2384855)
            if s:
                set_text_keep_format(s, f"{top30_sales_rounded} 백만 RMB", auto_fit=True)

            # TOP 30 점유율
            s = find_shape(slide3, 555287, 3209006)
            if s:
                set_text_keep_format(s, f"전체 매출의 {top30_sales_raw/gz_1q_total*100:.0f}% ", auto_fit=True)

            # 신규 고객사 수 → 공란
            s = find_shape(slide3, 554991, 3787841)
            if s:
                set_text_keep_format(s, "", auto_fit=True)

            # 사업계획 30대 고객사 +/- 항목 → 공란
            for left, top in [(468733, 5262651), (1612155, 5260628), (2963753, 5271738)]:
                s = find_shape(slide3, left, top)
                if s:
                    clear_all_text(s)

            # 내장 엑셀 개체 (OLE) → 삭제하고 TOP 30 표로 교체
            ole_shape = None
            for s in slide3.shapes:
                if s.shape_type == 7:  # EMBEDDED_OLE_OBJECT
                    ole_shape = s
                    break
            if ole_shape:
                ol, ot, ow, oh = ole_shape.left, ole_shape.top, ole_shape.width, ole_shape.height
                sp = ole_shape._element
                sp.getparent().remove(sp)

                # TOP 30 고객사 표 (원본 PPT와 동일 컬럼)
                top30_data = df_full_cust.head(30)
                n_rows = len(top30_data) + 4  # 헤더 + 30행 + 합계3행
                n_cols = 11
                max_h = 6858000 - ot - 50000
                tbl_shape = slide3.shapes.add_table(n_rows, n_cols, ol, ot, ow, max_h)
                tbl = tbl_shape.table

                col_widths = [370000, 550000, 450000, 650000, 650000, 680000, 680000, 550000, 620000, 480000, 550000]
                for ci, w in enumerate(col_widths):
                    tbl.columns[ci].width = w

                headers = ["순위", "고객사", "채널", "25년매출", "사업계획", "25년동기", "예상매출", "성장률", "달성률", "점유율", "전년대비\n순위"]
                for ci, h in enumerate(headers):
                    cell = tbl.cell(0, ci)
                    cell.text = h
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(5)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x20, 0x38, 0x64)
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                total_s = df_full_cust.head(30)["예상매출"].sum()
                for ri, (_, row) in enumerate(top30_data.iterrows()):
                    dr = ri + 1
                    tbl.cell(dr, 0).text = str(ri + 1)
                    tbl.cell(dr, 1).text = str(row["고객사"])
                    tbl.cell(dr, 2).text = str(row["채널"])
                    tbl.cell(dr, 3).text = f"{row['25년매출']/1000:.0f}" if row["25년매출"] else "-"
                    tbl.cell(dr, 4).text = f"{row['사업계획']/1000:.0f}" if row["사업계획"] else "-"
                    tbl.cell(dr, 5).text = f"{row['25년동기']/1000:.0f}" if row["25년동기"] else "-"
                    tbl.cell(dr, 6).text = f"{row['예상매출']/1000:.0f}" if row["예상매출"] else "-"
                    tbl.cell(dr, 7).text = f"{row['성장률']:.0%}" if row["성장률"] and row["25년동기"] else "-"
                    tbl.cell(dr, 8).text = f"{row['달성률']:.0%}" if row["달성률"] else "-"
                    tbl.cell(dr, 9).text = f"{row['점유율']:.0%}" if row["점유율"] else "-"
                    tbl.cell(dr, 10).text = str(row["전년대비순위"]) if row["전년대비순위"] else "-"

                    for ci in range(n_cols):
                        for p in tbl.cell(dr, ci).text_frame.paragraphs:
                            p.font.size = Pt(5)
                            p.alignment = PP_ALIGN.CENTER
                    if ri % 2 == 0:
                        for ci in range(n_cols):
                            tbl.cell(dr, ci).fill.solid()
                            tbl.cell(dr, ci).fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)

                # 하단 합계행 (TOP 10/20/30) - 반올림 합산
                for si, (label, cnt) in enumerate([("TOP 10 고객사 합계", 10), ("TOP 20 고객사 합계", 20), ("TOP 30 고객사 합계", 30)]):
                    sr = len(top30_data) + 1 + si
                    s_sum = sum(round(v / 1000) for v in df_full_cust.head(cnt)["예상매출"])
                    s_pct = df_full_cust.head(cnt)["예상매출"].sum() / gz_1q_total * 100 if gz_1q_total else 0
                    tbl.cell(sr, 0).text = ""
                    tbl.cell(sr, 1).text = label
                    for ci in range(2, 6):
                        tbl.cell(sr, ci).text = ""
                    tbl.cell(sr, 6).text = f"{s_sum}"
                    for ci in range(7, 9):
                        tbl.cell(sr, ci).text = ""
                    tbl.cell(sr, 9).text = f"{s_pct:.0f}%"
                    tbl.cell(sr, 10).text = ""
                    for ci in range(n_cols):
                        cell = tbl.cell(sr, ci)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0x20, 0x38, 0x64)
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(5)
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER
                            for r in p.runs:
                                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # ════════════════════════════════════
            # 슬라이드 4: 채널별 실적 분석
            # ════════════════════════════════════
            slide4 = prs.slides[3]

            # 전체 고객사 기반 채널별 데이터 (시트2에서 추출, PPT 원본과 일치)
            total_ch_sales = sum(channel_1q[ch]["sales"] for ch in channel_1q)

            # 채널별 위치 매핑
            channel_ui = [
                {"name": "온라인", "sales_pos": (909387, 1906415), "plan_pos": (2016744, 1897179),
                 "share_pos": (2738605, 1176094), "growth_pos": (1035628, 2222317),
                 "table_pos": (3062511, 1812408), "progress_pos": (1152766, 2904559)},
                {"name": "왕홍", "sales_pos": (6441323, 1859970), "plan_pos": (7530208, 1859970),
                 "share_pos": (10640781, 1175789), "growth_pos": (6576800, 2185108),
                 "table_pos": (8603879, 1774245), "progress_pos": (7950098, 2958627)},
                {"name": "오프라인", "sales_pos": (887880, 4578673), "plan_pos": (1995237, 4569437),
                 "share_pos": (2717098, 3848352), "growth_pos": (1009468, 4892526),
                 "table_pos": (3034471, 4528167), "progress_pos": (1133598, 5629787)},
                {"name": "수출", "sales_pos": (7493644, 4587737), "plan_pos": (8573296, 4587737),
                 "share_pos": (10667867, 3857376), "growth_pos": (7564469, 4903639),
                 "table_pos": (9601432, 4527366), "progress_pos": (7936878, 5637168)},
            ]

            for ch_info in channel_ui:
                ch_name = ch_info["name"]
                ch_data = channel_1q.get(ch_name, {"sales": 0, "plan": 0, "prev": 0})
                sales = ch_data["sales"]
                plan = ch_data["plan"]
                share = sales / total_ch_sales * 100 if total_ch_sales else 0
                plan_rate = sales / plan * 100 if plan else 0
                yoy = channel_yoy_ppt.get(ch_name, 0)

                # 예상 매출 값
                s = find_shape(slide4, *ch_info["sales_pos"])
                if s:
                    set_text_keep_format(s, f"{sales/1000:.0f}")

                # 계획 대비
                s = find_shape(slide4, *ch_info["plan_pos"])
                if s:
                    set_text_keep_format(s, f"{plan_rate:.0f}%")

                # 비중
                s = find_shape(slide4, *ch_info["share_pos"])
                if s:
                    set_text_keep_format(s, f"비중 {share:.0f}%")

                # 전년비
                s = find_shape(slide4, *ch_info["growth_pos"])
                if s:
                    if abs(yoy) < 0.5:
                        set_text_keep_format(s, " 전년비 = 0%")
                    else:
                        arrow = "▲" if yoy > 0 else "▼"
                        set_text_keep_format(s, f" 전년비 {arrow} {abs(yoy):.0f}%")
                    if s.text_frame.paragraphs[0].runs:
                        run = s.text_frame.paragraphs[0].runs[0]
                        if yoy > 0:
                            run.font.color.rgb = RGBColor(0x38, 0x57, 0x23)
                        elif yoy < 0:
                            run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)

                # TOP 3 테이블
                tbl_shape = find_shape(slide4, *ch_info["table_pos"])
                if tbl_shape and tbl_shape.has_table:
                    tbl = tbl_shape.table
                    cn_map = {"온라인": "线上", "왕홍": "红人", "오프라인": "线下", "수출": "出口"}
                    df_ch_top = df_full_cust[df_full_cust["채널"] == ch_name].head(3)
                    for ri, (_, row) in enumerate(df_ch_top.iterrows()):
                        if ri < len(list(tbl.rows)):
                            set_table_cell(tbl, ri, 0, str(row["고객사"]))
                            set_table_cell(tbl, ri, 2, f"{row['예상매출']/1000:.0f}")

                # 주요 진행사항 → 공란 (Text 19 여러 개)
                s = find_shape(slide4, *ch_info["progress_pos"])
                if s:
                    clear_all_text(s)
                # 추가 진행사항 텍스트도 공란
                extra_progress = {
                    "온라인": [(1137218, 5854383), (1663472, 3447188)],
                    "왕홍": [],
                    "오프라인": [(1137218, 5854383)],
                    "수출": [],
                }
                for ep_left, ep_top in extra_progress.get(ch_name, []):
                    ep_s = find_shape(slide4, ep_left, ep_top)
                    if ep_s:
                        clear_all_text(ep_s)

            # 3/6대비 값 → 공란
            for left, top in [(1467107, 1923148), (7000290, 1884318), (1477355, 4591829), (8066264, 4610670)]:
                s = find_shape(slide4, left, top)
                if s:
                    set_text_keep_format(s, "")

            # 가운데 파이 차트 → 이미지로 교체
            chart_s4 = None
            for s in slide4.shapes:
                if s.shape_type == 3:
                    chart_s4 = s
                    break
            if chart_s4:
                cl, ct, cw, ch_ = chart_s4.left, chart_s4.top, chart_s4.width, chart_s4.height
                sp = chart_s4._element
                sp.getparent().remove(sp)

                ch_pie_data = pd.DataFrame([
                    {"채널": ch, "매출": channel_1q[ch]["sales"]}
                    for ch in ["온라인", "왕홍", "오프라인", "수출"]
                ])
                ppt_ch_colors = {"온라인": "#223862", "왕홍": "#BE0201", "오프라인": "#555555", "수출": "#D1D6E7"}
                fig_pie = px.pie(ch_pie_data, names="채널", values="매출",
                                 color="채널", color_discrete_map=ppt_ch_colors, hole=0.4)
                fig_pie.update_traces(textposition="inside", textinfo="label+percent", textfont_size=13)
                fig_pie.update_layout(
                    title=dict(text="채널별 매출 비중", font=dict(size=14, color="#1F2937"), x=0.5, y=0.55),
                    showlegend=True,
                    legend=dict(orientation="h", yanchor="top", y=-0.02, xanchor="center", x=0.5, font=dict(size=10)),
                    margin=dict(l=5, r=5, t=5, b=40), font=dict(size=12),
                )
                # 타원(2968468x2968468) 안에 맞추기: 타원 위치 기준으로 센터링
                oval_left, oval_top, oval_size = 4612920, 2418073, 2968468
                slide4.shapes.add_picture(fig_to_image(fig_pie, 500, 500),
                    Emu(oval_left), Emu(oval_top), Emu(oval_size), Emu(oval_size))

            # ════════════════════════════════════
            # 슬라이드 5: 신규 성장 프로젝트 → 양식 유지, 내용 공란
            # ════════════════════════════════════
            slide5 = prs.slides[4]
            # 프로젝트명, 진행사항, 우측 카테고리 내용 텍스트 공란 처리
            # 좌측 프로젝트명 (TextBox 84,18,85,86,87,88,89,90)
            proj_name_positions = [
                (681336, 1941105), (681336, 2457226), (692024, 2981776),
                (681336, 3511074), (681336, 4045449), (681335, 4563277),
                (681336, 5089531), (681335, 5606803),
            ]
            for left, top in proj_name_positions:
                s = find_shape(slide5, left, top)
                if s:
                    clear_all_text(s)

            # 좌측 진행사항 (TextBox 5,6,8,11,23,25,24,26)
            prog_positions = [
                (2424030, 1943974), (2424030, 2460725), (2424029, 2977476),
                (2420649, 3510079), (2420647, 4055330), (2420646, 4564442),
                (2420646, 5081531), (2420646, 5644553),
            ]
            for left, top in prog_positions:
                s = find_shape(slide5, left, top)
                if s:
                    clear_all_text(s)

            # 우측 카테고리 내용 (TextBox 43,45,47,52)
            right_positions = [
                (7651391, 1613726), (7651391, 2858841),
                (7651391, 4073453), (7651391, 5357902),
            ]
            for left, top in right_positions:
                s = find_shape(slide5, left, top)
                if s:
                    clear_all_text(s)

            # ── PPT 저장 ──
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            try:
                os.unlink(tmp.name)
            except Exception:
                pass

        st.success("PPT 생성 완료!")
        st.download_button(
            label="PPT 다운로드",
            data=buffer,
            file_name=f"코스맥스_광저우_보고_{report_date.strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

# ═══════════════════════════════════════
# 공통 하단 푸터
# ═══════════════════════════════════════
st.markdown("---")
st.markdown(f"""
<div style="text-align:center;padding:16px 0 8px 0;color:#9CA3AF;font-size:0.8rem;">
    Last updated: {datetime.now().strftime("%Y-%m-%d")} | Created by <b style="color:#6B7280;">[SUJUNG CHOI]</b>
</div>
""", unsafe_allow_html=True)
